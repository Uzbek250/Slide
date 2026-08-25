"""
Screenshot qilingan PNG rasmlarni bitta .pptx faylga yig'adi.
Har bir rasm — bitta slaydga to'liq (full-bleed) joylashtiriladi.
"""
from pptx import Presentation
from pptx.util import Emu

# 1920x1080 px, 2x device_scale_factor bilan chiqarilgan rasm uchun
# 16:9 pptx slayd o'lchami (standart: 13.333 x 7.5 inch)
SLIDE_WIDTH_EMU = Emu(12192000)   # 13.333 in
SLIDE_HEIGHT_EMU = Emu(6858000)   # 7.5 in


def build_pptx(image_paths: list[str], output_path: str, deck_title: str = "Presentation") -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH_EMU
    prs.slide_height = SLIDE_HEIGHT_EMU

    blank_layout = prs.slide_layouts[6]  # to'liq bo'sh layout

    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_path, left=0, top=0,
            width=SLIDE_WIDTH_EMU, height=SLIDE_HEIGHT_EMU
        )

    prs.core_properties.title = deck_title
    prs.save(output_path)
    return output_path
